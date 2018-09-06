# coding: utf-8

import os
import re
import shutil
import glob
import pptx
from PIL import Image

INPUT_DIR = r"C:\Users\FU25166\Documents\Pythoncodes\20180906_File search\testdata"
OUTPUT_DIR = r"C:\Users\FU25166\Documents\Pythoncodes\20180906_File search\output"

# 1次列目をまとめる
word = input("検索する言葉（サンプル名）を入力してください：")
print(word)

makefolder = OUTPUT_DIR + "\\" + word
os.makedirs(makefolder)
for foldername, subfolders, filenames in os.walk(INPUT_DIR):
    for filename in filenames:
        mojiretsu = foldername + "\\" + filename
        kensaku1 = re.compile(word)
        mo1 = kensaku1.search(mojiretsu)
        if mo1 == None:
            continue
        else:
            shutil.copy(mojiretsu, makefolder)

# 2次列目をまとめる
word2 = str(input("検索する言葉（場所）を入力してください："))
makefolder2 = OUTPUT_DIR + "\\" + word2
os.makedirs(makefolder2)
for foldername, subfolders, filenames in os.walk(makefolder):
    for filename in filenames:
        mojiretsu2 = foldername + "\\" + filename
        kensaku2 = re.compile(word2)
        mo2 = kensaku2.search(mojiretsu2)
        if mo2 == None:
            continue
        else:
            shutil.copy(mojiretsu2, makefolder2)


# パワポに展開
path_list = glob.glob(r"C:\Users\FU25166\Documents\Pythoncodes\20180906_File search\testdata\*.JPG")  # .JPGとつくファイルをリストで取得
print(path_list)

ppt = pptx.Presentation()
width = ppt.slide_width
height = ppt.slide_height
blank_slide_layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(blank_slide_layout)

for j, i in enumerate(path_list):
    img = Image.open(i)
    img_resize = img.resize((256, 256))
    img_resize.save(i)
    pic = slide.shapes.add_picture(i, 0, 0)

    pic.left = int((width - pic.width) / (j + 1))
    pic.top = int((height - pic.height) / (1000))

ppt.save('figure.pptx')

# #パワポに展開
# import glob
# import pptx
# from PIL import Image
#
# path_list=glob.glob("C:\\Users\\VL58003\\Desktop\\Place\\"+word2+"\\*.jpg") #.jpgとつくファイルをリストで取得
#
# ppt=pptx.Presentation()
# width=ppt.slide_width
# height=ppt.slide_height
# blank_slide_layout=ppt.slide_layouts[6]
# slide=ppt.slides.add_slide(blank_slide_layout)
#
# for j,i in enumerate(path_list):
#     img=Image.open(i)
#     img_resize = img.resize((128, 128))
#     img_resize.save(i)
#     pic=slide.shapes.add_picture(i,0,0)
#
#     pic.left = int(j*2000000)
#     pic.top  = int( ( height - pic.height ) / (1000) )
#
# ppt.save('figure.pptx')
